"""Flask app module."""

from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path
from uuid import uuid4
import json
import mimetypes
import os
import shutil
from typing import Any

from flask import Flask, abort, flash, jsonify, redirect, render_template, request, send_file, url_for
from werkzeug.utils import secure_filename

from studyscribe.core import config, db
from studyscribe.services.audio import save_audio
from studyscribe.services.export import build_session_export
from studyscribe.services.gemini import GeminiError, answer_question, generate_notes
from studyscribe.services.jobs import create_job, enqueue_job, get_job
from studyscribe.services.transcribe import TranscriptionError, load_transcript, transcribe_audio
from studyscribe.services.retrieval import build_chunks, retrieve_chunks


BASE_DIR = Path(__file__).resolve().parent
app = Flask(
    __name__,
    static_folder=str(BASE_DIR / "web" / "static"),
    template_folder=str(BASE_DIR / "web" / "templates"),
)
app.secret_key = os.getenv("FLASK_SECRET", "studyscribe-dev")

ALLOWED_AUDIO_EXTENSIONS = {".wav", ".mp3", ".m4a", ".aac", ".flac", ".ogg"}
ALLOWED_ATTACHMENT_EXTENSIONS = {".pdf", ".ppt", ".pptx", ".doc", ".docx"}
ALLOWED_ATTACHMENT_MIME_TYPES = {
    "application/pdf",
    "application/msword",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.ms-powerpoint",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
SEGMENT_TAGS = {"IMPORTANT", "CONFUSING", "EXAM-SIGNAL"}


@app.context_processor
def inject_config():
    return {"config": config.settings}


def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def _module_dir(module_id: str) -> Path:
    return config.DATA_DIR / "modules" / module_id


def _session_dir(module_id: str, session_id: str) -> Path:
    return _module_dir(module_id) / "sessions" / session_id


def _ensure_session_dirs(session_dir: Path) -> None:
    for name in ("audio", "attachments", "transcript", "notes", "exports", "work"):
        (session_dir / name).mkdir(parents=True, exist_ok=True)


def _annotations_path(session_dir: Path) -> Path:
    return session_dir / "annotations.json"


def _load_annotations(session_dir: Path) -> dict[str, Any]:
    path = _annotations_path(session_dir)
    if not path.exists():
        return {
            "tags": {},
            "notes": "",
            "notes_html": "",
            "notes_markdown": "",
            "session_tags": [],
        }
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return {
            "tags": {},
            "notes": "",
            "notes_html": "",
            "notes_markdown": "",
            "session_tags": [],
        }


def _save_annotations(session_dir: Path, payload: dict[str, Any]) -> None:
    path = _annotations_path(session_dir)
    path.write_text(json.dumps(payload, indent=2), encoding="utf-8")


def _load_ai_notes(session_dir: Path) -> tuple[str, list[str]]:
    notes_path = session_dir / "notes" / "ai_notes.md"
    notes_json = session_dir / "notes" / "ai_notes.json"
    notes = notes_path.read_text(encoding="utf-8") if notes_path.exists() else ""
    suggested_tags: list[str] = []
    if notes_json.exists():
        try:
            data = json.loads(notes_json.read_text(encoding="utf-8"))
            suggested_tags = data.get("suggested_tags") or []
        except json.JSONDecodeError:
            suggested_tags = []
    return notes, suggested_tags


def _load_extracted_sources(session_dir: Path) -> list[dict]:
    sources_path = session_dir / "attachments" / "extracted_sources.json"
    if not sources_path.exists():
        return []
    try:
        return json.loads(sources_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return []


def _extract_pdf_text(path: Path) -> tuple[str, list[dict]]:
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        return "", []
    text_blocks = []
    sources = []
    with pdfplumber.open(path) as pdf:
        for page_index, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if text:
                text_blocks.append(text)
                sources.append(
                    {
                        "source_id": f"att_{path.name}_p{page_index}",
                        "kind": "attachment",
                        "file_name": path.name,
                        "mime": "application/pdf",
                        "page": page_index,
                        "text": text,
                    }
                )
    return "\n\n".join(text_blocks), sources


def _extract_docx_text(path: Path) -> tuple[str, list[dict]]:
    try:
        import docx  # type: ignore
    except ImportError:
        return "", []
    document = docx.Document(str(path))
    paragraphs = [para.text.strip() for para in document.paragraphs if para.text.strip()]
    text = "\n".join(paragraphs)
    sources = []
    if text:
        sources.append(
            {
                "source_id": f"att_{path.name}",
                "kind": "attachment",
                "file_name": path.name,
                "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "page": None,
                "text": text,
            }
        )
    return text, sources


def _extract_pptx_text(path: Path) -> tuple[str, list[dict]]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return "", []
    presentation = Presentation(str(path))
    text_blocks = []
    sources = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content = str(shape.text).strip()
                if content:
                    slide_text.append(content)
        text = "\n".join(slide_text).strip()
        if text:
            text_blocks.append(text)
            sources.append(
                {
                    "source_id": f"att_{path.name}_s{slide_index}",
                    "kind": "attachment",
                    "file_name": path.name,
                    "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    "slide": slide_index,
                    "text": text,
                }
            )
    return "\n\n".join(text_blocks), sources


def _rebuild_attachment_index(session_dir: Path) -> None:
    attachments_dir = session_dir / "attachments"
    extracted_text: list[str] = []
    sources: list[dict] = []
    for path in attachments_dir.iterdir():
        if not path.is_file():
            continue
        ext = path.suffix.lower()
        if ext not in ALLOWED_ATTACHMENT_EXTENSIONS:
            continue
        if ext == ".pdf":
            text, extracted_sources = _extract_pdf_text(path)
        elif ext in {".doc", ".docx"}:
            text, extracted_sources = _extract_docx_text(path)
        elif ext in {".ppt", ".pptx"}:
            text, extracted_sources = _extract_pptx_text(path)
        else:
            text, extracted_sources = "", []
        if text:
            extracted_text.append(text)
        sources.extend(extracted_sources)
    (attachments_dir / "extracted.txt").write_text("\n\n".join(extracted_text), encoding="utf-8")
    (attachments_dir / "extracted_sources.json").write_text(
        json.dumps(sources, indent=2), encoding="utf-8"
    )


def _collect_attachment_files(session_dir: Path) -> list[dict]:
    attachments_dir = session_dir / "attachments"
    if not attachments_dir.exists():
        return []
    extracted_sources = _load_extracted_sources(session_dir)
    files = []
    for path in sorted(attachments_dir.iterdir()):
        if not path.is_file():
            continue
        if path.name in {"extracted.txt", "extracted_sources.json"}:
            continue
        files.append(
            {
                "name": path.name,
                "size": _format_size(path.stat().st_size),
                "has_text": any(source.get("file_name") == path.name for source in extracted_sources),
            }
        )
    return files


def _format_size(size_bytes: int) -> str:
    if size_bytes < 1024:
        return f"{size_bytes} B"
    if size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    return f"{size_bytes / (1024 * 1024):.1f} MB"


def _collect_files(directory: Path, allowed_extensions: set[str] | None = None) -> list[dict]:
    if not directory.exists():
        return []
    files = []
    for path in sorted(directory.iterdir(), key=lambda item: item.stat().st_mtime, reverse=True):
        if not path.is_file():
            continue
        if allowed_extensions is not None and path.suffix.lower() not in allowed_extensions:
            continue
        files.append(
            {
                "name": path.name,
                "size": _format_size(path.stat().st_size),
            }
        )
    return files


def _collect_audio_files(session_dir: Path) -> list[dict]:
    return _collect_files(session_dir / "audio", ALLOWED_AUDIO_EXTENSIONS)


def _select_latest_audio(session_dir: Path) -> Path | None:
    audio_dir = session_dir / "audio"
    if not audio_dir.exists():
        return None
    files = [p for p in audio_dir.iterdir() if p.is_file()]
    if not files:
        return None
    return max(files, key=lambda p: p.stat().st_mtime)


def _clear_transcript(session_dir: Path) -> None:
    transcript_dir = session_dir / "transcript"
    if transcript_dir.exists():
        shutil.rmtree(transcript_dir)
    annotations = _load_annotations(session_dir)
    annotations["tags"] = {}
    _save_annotations(session_dir, annotations)


def _format_ts(seconds: float) -> str:
    minutes = int(seconds // 60)
    secs = int(seconds % 60)
    return f"{minutes:02d}:{secs:02d}"


def _format_datetime(value: str | None) -> str:
    if not value:
        return ""
    try:
        parsed = datetime.fromisoformat(value.replace("Z", "+00:00"))
    except ValueError:
        return value
    return parsed.astimezone(timezone.utc).strftime("%d %b %Y, %I:%M %p")


def _wants_json() -> bool:
    accept = request.headers.get("Accept", "")
    return "application/json" in accept.lower()


def _json_error(message: str, status: int = 400):
    return jsonify({"ok": False, "error": message}), status


def _json_not_implemented(message: str = "Not implemented in Sprint 2 (UI/UX parity)."):
    return _json_error(message, status=501)


def _store_ai_message(session_id: str, role: str, content: str) -> int:
    return db.execute_returning_id(
        "INSERT INTO ai_messages (session_id, role, content, created_at) VALUES (?, ?, ?, ?)",
        (session_id, role, content, _now_iso()),
    )


def _store_ai_sources(message_id: int, sources: list[dict], session_name: str) -> None:
    for source in sources:
        db.execute(
            """
            INSERT INTO ai_message_sources
                (message_id, source_id, kind, label, snippet, session_name, url, source_json)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                message_id,
                source.get("source_id") or str(source.get("id")),
                source.get("kind", "transcript"),
                source.get("title") or source.get("label") or "Source",
                source.get("excerpt"),
                session_name,
                source.get("open_url"),
                json.dumps(source),
            ),
        )


def _load_ai_messages(session_id: str) -> list[dict]:
    messages = db.fetch_all(
        "SELECT * FROM ai_messages WHERE session_id = ? ORDER BY id ASC", (session_id,)
    )
    result = []
    for message in messages:
        sources = db.fetch_all(
            "SELECT * FROM ai_message_sources WHERE message_id = ? ORDER BY id ASC",
            (message["id"],),
        )
        mapped_sources = []
        for source in sources:
            payload = dict(source)
            if "url" in payload:
                payload["open_url"] = payload.pop("url")
            mapped_sources.append(payload)
        result.append(
            {
                "id": message["id"],
                "session_id": message["session_id"],
                "role": message["role"],
                "content": message["content"],
                "created_at": message["created_at"],
                "sources": mapped_sources,
            }
        )
    return result


def _build_notes_prompt(transcript: list[dict], attachment_sources: list[dict]) -> str:
    transcript_text = "\n".join(
        f"[{seg.get('start'):.2f}-{seg.get('end'):.2f}] {seg.get('text', '')}"
        for seg in transcript
    )
    attachments_text = "\n\n".join(source.get("text", "") for source in attachment_sources)
    return (
        "You are StudyScribe. Create structured study notes from the content below.\n"
        "Return JSON with keys: summary (string), suggested_tags (array of strings), notes_markdown (string).\n\n"
        "Transcript:\n"
        f"{transcript_text}\n\n"
        "Attachments:\n"
        f"{attachments_text}\n"
    )


def _build_qa_prompt(question: str, context: str) -> str:
    return (
        "Answer the question using ONLY the provided context. "
        "Return JSON with keys: answer (string) and answer_markdown (string).\n\n"
        f"Question:\n{question}\n\n"
        f"Context:\n{context}\n"
    )


def _handle_qa_request(session_id: str, question: str, scope: str):
    session = db.fetch_one("SELECT * FROM sessions WHERE id = ?", (session_id,))
    if not session:
        return {"error": "Session not found."}, 404
    session_rows = [session]
    if scope == "module":
        session_rows = db.fetch_all(
            "SELECT * FROM sessions WHERE module_id = ? ORDER BY created_at DESC",
            (session["module_id"],),
        )

    all_chunks: list[dict] = []
    all_attachment_sources: list[dict] = []
    for row in session_rows:
        session_dir = _session_dir(row["module_id"], row["id"])
        transcript = load_transcript(session_dir / "transcript" / "transcript.json")
        chunks_path = session_dir / "transcript" / "chunks.json"
        if chunks_path.exists():
            try:
                chunks = json.loads(chunks_path.read_text(encoding="utf-8"))
            except json.JSONDecodeError:
                chunks = build_chunks(transcript)
        else:
            chunks = build_chunks(transcript)
        for chunk in chunks:
            chunk["session_id"] = row["id"]
            chunk["session_name"] = row["name"]
        all_chunks.extend(chunks)
        attachment_sources = _load_extracted_sources(session_dir)
        for source in attachment_sources:
            source["session_id"] = row["id"]
            source["session_name"] = row["name"]
        all_attachment_sources.extend(attachment_sources)

    if not all_chunks and not all_attachment_sources:
        return {"error": "Upload transcript or attachments to enable Q&A."}, 400

    transcript_hits = retrieve_chunks(question, all_chunks, k=6)
    sources: list[dict] = []
    for idx, chunk in enumerate(transcript_hits, start=1):
        start = float(chunk.get("start", 0.0))
        end = float(chunk.get("end", 0.0))
        segment_id = chunk.get("segment_ids", [None])[0]
        chunk_session_id = chunk.get("session_id", session_id)
        sources.append(
            {
                "id": idx,
                "source_id": f"src_{idx}",
                "kind": "transcript",
                "title": f"Transcript [{_format_ts(start)}–{_format_ts(end)}]",
                "excerpt": (chunk.get("text", "") or "")[:220],
                "locator": {
                    "type": "transcript",
                    "session_id": chunk_session_id,
                    "segment_id": segment_id,
                    "t_start": start,
                    "t_end": end,
                    "t_start_ms": int(start * 1000),
                    "t_end_ms": int(end * 1000),
                    "anchor": f"seg-{segment_id}",
                },
            }
        )

    attachment_chunks = [
        {"text": source.get("text", ""), "source": source} for source in all_attachment_sources
    ]
    attachment_hits = retrieve_chunks(question, attachment_chunks, k=3)
    for hit in attachment_hits:
        source = hit.get("source", {})
        file_name = source.get("file_name")
        page = source.get("page") or source.get("slide")
        open_url = url_for(
            "open_attachment",
            module_id=session["module_id"],
            session_id=source.get("session_id", session_id),
            filename=file_name,
        )
        if page:
            open_url = f"{open_url}#page={page}"
        sources.append(
            {
                "id": len(sources) + 1,
                "source_id": source.get("source_id") or f"src_att_{len(sources) + 1}",
                "kind": "attachment",
                "title": f"Attachment: {file_name}",
                "excerpt": (source.get("text") or "")[:220],
                "locator": {
                    "type": "attachment",
                    "attachment_id": file_name,
                    "file_name": file_name,
                    "mime": source.get("mime"),
                    "page": source.get("page"),
                    "chunk_id": source.get("source_id"),
                },
                "open_url": open_url,
            }
        )

    context = "\n\n".join(f"[{source['id']}] {source['excerpt']}" for source in sources if source.get("excerpt"))
    prompt = _build_qa_prompt(question, context)

    try:
        answer = answer_question(prompt, sources)
    except GeminiError as exc:
        return {"error": exc.user_message}, 500

    user_message_id = _store_ai_message(session_id, "user", question)
    assistant_message_id = _store_ai_message(session_id, "assistant", answer.answer_markdown)
    _store_ai_sources(assistant_message_id, sources, session["name"])

    notes_dir = session_dir / "notes"
    notes_dir.mkdir(parents=True, exist_ok=True)
    (notes_dir / "last_answer.json").write_text(
        json.dumps(
            {
                "answer": answer.answer,
                "answer_markdown": answer.answer_markdown,
                "sources": sources,
                "context_sources": sources,
                "scope": scope,
                "question": question,
            },
            indent=2,
        ),
        encoding="utf-8",
    )

    return {
        "answer": answer.answer,
        "answer_markdown": answer.answer_markdown,
        "sources": sources,
        "user_message_id": user_message_id,
        "assistant_message_id": assistant_message_id,
    }

def _init() -> None:
    config.DATA_DIR.mkdir(parents=True, exist_ok=True)
    db.init_db()


@app.template_filter("format_ts")
def _format_ts_filter(seconds: float) -> str:
    return _format_ts(seconds)


@app.template_filter("format_dt")
def _format_dt_filter(value: str | None) -> str:
    return _format_datetime(value)


@app.route("/")
def index():
    return redirect(url_for("home"))


@app.route("/home")
def home():
    modules = db.fetch_all("SELECT * FROM modules ORDER BY created_at DESC")
    return render_template("index.html", modules=[dict(row) for row in modules])


@app.route("/modules", methods=["POST"])
def create_module():
    name = (request.form.get("name") or "").strip()
    if not name:
        flash("Module name is required.", "error")
        return redirect(url_for("home")), 400
    module_id = str(uuid4())
    db.execute(
        "INSERT INTO modules (id, name, created_at) VALUES (?, ?, ?)",
        (module_id, name, _now_iso()),
    )
    _module_dir(module_id).mkdir(parents=True, exist_ok=True)
    flash("Module created.", "success")
    return redirect(url_for("view_module", module_id=module_id))


@app.route("/modules/<module_id>")
def view_module(module_id: str):
    module = db.fetch_one("SELECT * FROM modules WHERE id = ?", (module_id,))
    if not module:
        abort(404)
    sessions = db.fetch_all(
        "SELECT * FROM sessions WHERE module_id = ? ORDER BY created_at DESC", (module_id,)
    )
    modules = db.fetch_all("SELECT * FROM modules ORDER BY created_at DESC")
    return render_template(
        "module.html",
        module=dict(module),
        sessions=[dict(row) for row in sessions],
        modules=[dict(row) for row in modules],
    )


@app.route("/modules/<module_id>/sessions", methods=["POST"])
def create_session(module_id: str):
    module = db.fetch_one("SELECT * FROM modules WHERE id = ?", (module_id,))
    if not module:
        abort(404)
    name = (request.form.get("name") or "").strip() or "Untitled"
    session_id = str(uuid4())
    db.execute(
        "INSERT INTO sessions (id, module_id, name, created_at) VALUES (?, ?, ?, ?)",
        (session_id, module_id, name, _now_iso()),
    )
    session_dir = _session_dir(module_id, session_id)
    session_dir.mkdir(parents=True, exist_ok=True)
    _ensure_session_dirs(session_dir)
    rename = "1" if name == "Untitled" else None
    return redirect(
        url_for("view_session", session_id=session_id, rename=rename) if rename else url_for("view_session", session_id=session_id)
    )


@app.route("/modules/<module_id>/sessions/<session_id>")
def legacy_view_session(module_id: str, session_id: str):
    return redirect(url_for("view_session", session_id=session_id))


@app.route("/sessions/<session_id>")
def view_session(session_id: str):
    session = db.fetch_one("SELECT * FROM sessions WHERE id = ?", (session_id,))
    if not session:
        abort(404)
    module = db.fetch_one("SELECT * FROM modules WHERE id = ?", (session["module_id"],))
    if not module:
        abort(404)
    sessions = db.fetch_all(
        "SELECT * FROM sessions WHERE module_id = ? ORDER BY created_at DESC",
        (session["module_id"],),
    )
    session_dir = _session_dir(session["module_id"], session_id)
    _ensure_session_dirs(session_dir)
    transcript_path = session_dir / "transcript" / "transcript.json"
    transcript = load_transcript(transcript_path)
    job_id = request.args.get("job_id")
    modules = db.fetch_all("SELECT * FROM modules ORDER BY created_at DESC")
    audio_files = _collect_audio_files(session_dir)
    attachment_files = _collect_attachment_files(session_dir)
    attachments_with_text = {
        item["name"] for item in attachment_files if item.get("has_text")
    }
    attachment_warning = ""
    annotations = _load_annotations(session_dir)
    notes_html = annotations.get("notes_html", "")
    notes_markdown = annotations.get("notes_markdown", "")
    ai_notes, suggested_tags = _load_ai_notes(session_dir)
    session_tags = annotations.get("session_tags", [])
    has_transcript = bool(transcript)
    has_attachments = bool(attachment_files)
    has_attachment_text = len(attachments_with_text) > 0
    has_notes = bool(annotations.get("notes") or notes_html or notes_markdown)
    has_generate_content = has_transcript or has_attachment_text or has_notes
    has_qa_content = has_transcript or has_attachment_text
    generate_hint = "Upload audio or attachments to generate notes."
    qa_hint = "Upload transcript or attachments to enable Q&A."
    if has_attachments and not has_attachment_text:
        attachment_warning = "Attachments uploaded, but no text could be extracted."
    for segment in transcript:
        segment_key = f"seg_{segment.get('segment_id')}"
        segment["tags"] = annotations.get("tags", {}).get(segment_key, [])
    session_meta = {
        "moduleId": session["module_id"],
        "sessionId": session_id,
        "moduleName": module["name"],
        "sessionName": session["name"],
        "autoRename": request.args.get("rename") == "1",
        "exportUrl": url_for("export_pack", module_id=module["id"], session_id=session_id),
        "generateUrl": url_for("start_notes", module_id=module["id"], session_id=session_id),
        "notesUrl": url_for("fetch_ai_notes", module_id=module["id"], session_id=session_id),
        "transcriptUrl": url_for("fetch_transcript", module_id=module["id"], session_id=session_id),
        "deleteTranscriptUrl": url_for("delete_transcript", module_id=module["id"], session_id=session_id),
        "segmentTagsUrl": url_for("update_segment_tags", module_id=module["id"], session_id=session_id),
        "qaAskUrl": url_for("api_ai_ask"),
        "qaMessagesUrl": url_for("api_ai_messages", session_id=session_id),
        "sourcePreviewUrl": url_for("api_source_preview"),
        "hasAudio": len(audio_files) > 0,
        "hasTranscript": has_transcript,
        "hasAttachments": has_attachments,
        "hasAttachmentText": has_attachment_text,
        "hasNotes": has_notes,
        "hasGenerateContent": has_generate_content,
        "hasQaContent": has_qa_content,
        "suggestedTags": suggested_tags,
    }
    return render_template(
        "session.html",
        module=dict(module),
        session=dict(session),
        sessions=[dict(row) for row in sessions],
        modules=[dict(row) for row in modules],
        audio_files=audio_files,
        attachment_files=attachment_files,
        attachments_with_text=attachments_with_text,
        attachment_warning=attachment_warning,
        transcript=transcript,
        annotations=annotations,
        notes_html=notes_html,
        notes_markdown=notes_markdown,
        ai_notes=ai_notes,
        suggested_tags=suggested_tags,
        session_tags=session_tags,
        has_generate_content=has_generate_content,
        has_qa_content=has_qa_content,
        generate_hint=generate_hint,
        qa_hint=qa_hint,
        session_meta=session_meta,
        job_id=job_id,
        transcript_url=url_for("fetch_transcript", module_id=module["id"], session_id=session_id),
    )


@app.route("/modules/<module_id>/sessions/<session_id>/upload-audio", methods=["POST"])
def upload_audio(module_id: str, session_id: str):
    session = db.fetch_one(
        "SELECT * FROM sessions WHERE id = ? AND module_id = ?", (session_id, module_id)
    )
    if not session:
        abort(404)
    file_storage = request.files.get("audio")
    if not file_storage or not file_storage.filename:
        if _wants_json():
            return _json_error("No audio file selected.", status=400)
        flash("No audio file selected.", "error")
        return redirect(url_for("view_session", session_id=session_id)), 400
    ext = Path(file_storage.filename).suffix.lower()
    if ext not in ALLOWED_AUDIO_EXTENSIONS:
        if _wants_json():
            return _json_error("Unsupported audio file type.", status=400)
        flash("Unsupported audio file type.", "error")
        return redirect(url_for("view_session", session_id=session_id)), 400
    session_dir = _session_dir(module_id, session_id)
    _ensure_session_dirs(session_dir)
    audio_dir = session_dir / "audio"
    existing_audio = list(audio_dir.glob("*")) if audio_dir.exists() else []
    replace = request.form.get("replace") == "1"
    if existing_audio and not replace:
        if _wants_json():
            return _json_error("Audio already uploaded. Replace the audio to continue.", status=400)
        flash("Audio already uploaded. Replace the audio to continue.", "error")
        return redirect(url_for("view_session", session_id=session_id)), 400
    if existing_audio and replace:
        shutil.rmtree(audio_dir)
        _clear_transcript(session_dir)
    saved_path = save_audio(file_storage, session_dir)
    if _wants_json():
        return jsonify({"ok": True, "filename": saved_path.name})
    flash(f"Audio saved to {saved_path.name}.", "success")
    return redirect(url_for("view_session", session_id=session_id))


@app.route("/modules/<module_id>/sessions/<session_id>/transcribe", methods=["POST"])
def start_transcription(module_id: str, session_id: str):
    session = db.fetch_one(
        "SELECT * FROM sessions WHERE id = ? AND module_id = ?", (session_id, module_id)
    )
    if not session:
        abort(404)
    session_dir = _session_dir(module_id, session_id)
    audio_path = _select_latest_audio(session_dir)
    if not audio_path:
        flash("Upload an audio file before transcribing.", "error")
        return redirect(url_for("view_session", session_id=session_id)), 400
    job_id = create_job("Queued for transcription.")
    try:
        enqueue_job(job_id, transcribe_audio, audio_path, session_dir)
    except TranscriptionError as exc:
        flash(exc.user_message, "error")
        return redirect(url_for("view_session", session_id=session_id)), 500
    flash("Transcription started.", "success")
    return redirect(url_for("view_session", session_id=session_id, job_id=job_id))


@app.route("/jobs/<job_id>")
def job_status(job_id: str):
    job = get_job(job_id)
    if not job:
        abort(404)
    return jsonify(
        {
            "id": job["id"],
            "status": job["status"],
            "progress": job["progress"],
            "message": job["message"],
            "result": job.get("result_path"),
        }
    )


@app.route("/modules/<module_id>/sessions/<session_id>/transcript")
def fetch_transcript(module_id: str, session_id: str):
    session = db.fetch_one(
        "SELECT * FROM sessions WHERE id = ? AND module_id = ?", (session_id, module_id)
    )
    if not session:
        abort(404)
    session_dir = _session_dir(module_id, session_id)
    transcript_path = session_dir / "transcript" / "transcript.json"
    transcript = load_transcript(transcript_path)
    annotations = _load_annotations(session_dir)
    for segment in transcript:
        segment_key = f"seg_{segment.get('segment_id')}"
        segment["tags"] = annotations.get("tags", {}).get(segment_key, [])
    html = render_template("_transcript_panel.html", transcript=transcript, annotations=annotations)
    return jsonify({"html": html, "has_transcript": bool(transcript)})


@app.route("/modules/<module_id>/sessions/<session_id>/upload-attachment", methods=["POST"])
def upload_attachment(module_id: str, session_id: str):
    session = db.fetch_one(
        "SELECT * FROM sessions WHERE id = ? AND module_id = ?", (session_id, module_id)
    )
    if not session:
        abort(404)
    files = request.files.getlist("attachment")
    if not files:
        if _wants_json():
            return _json_error("No attachment selected.", status=400)
        flash("No attachment selected.", "error")
        return redirect(url_for("view_session", session_id=session_id)), 400
    session_dir = _session_dir(module_id, session_id)
    _ensure_session_dirs(session_dir)
    attachments_dir = session_dir / "attachments"
    attachments_dir.mkdir(parents=True, exist_ok=True)
    saved_any = False
    pptx_uploaded = False
    for file_storage in files:
        if not file_storage or not file_storage.filename:
            continue
        filename = secure_filename(file_storage.filename or "") or "attachment"
        ext = Path(filename).suffix.lower()
        if ext not in ALLOWED_ATTACHMENT_EXTENSIONS:
            continue
        mime = file_storage.mimetype
        if mime and mime not in ALLOWED_ATTACHMENT_MIME_TYPES:
            continue
        if ext in {".ppt", ".pptx"}:
            pptx_uploaded = True
        dest = attachments_dir / filename
        file_storage.save(dest)
        saved_any = True
    if not saved_any:
        if _wants_json():
            return _json_error("Unsupported attachment type.", status=400)
        flash("Unsupported attachment type.", "error")
        return redirect(url_for("view_session", session_id=session_id)), 400
    _rebuild_attachment_index(session_dir)
    if pptx_uploaded:
        try:
            import pptx  # noqa: F401
        except ImportError:
            flash("python-pptx not installed; PPTX text extraction skipped.", "warning")
    if _wants_json():
        return jsonify({"ok": True})
    flash("Attachment uploaded.", "success")
    return redirect(url_for("view_session", session_id=session_id))


@app.route("/modules/<module_id>/sessions/<session_id>/delete-audio", methods=["POST"])
def delete_audio(module_id: str, session_id: str):
    session = db.fetch_one(
        "SELECT * FROM sessions WHERE id = ? AND module_id = ?", (session_id, module_id)
    )
    if not session:
        abort(404)
    filename = request.form.get("filename") or ""
    if not filename:
        if _wants_json():
            return _json_error("Filename is required.", status=400)
        flash("Filename is required.", "error")
        return redirect(url_for("view_session", session_id=session_id)), 400
    audio_path = _session_dir(module_id, session_id) / "audio" / filename
    if audio_path.exists():
        audio_path.unlink()
    if _wants_json():
        return jsonify({"ok": True})
    flash("Audio deleted.", "success")
    return redirect(url_for("view_session", session_id=session_id))


@app.route("/modules/<module_id>/sessions/<session_id>/delete-attachment", methods=["POST"])
def delete_attachment(module_id: str, session_id: str):
    session = db.fetch_one(
        "SELECT * FROM sessions WHERE id = ? AND module_id = ?", (session_id, module_id)
    )
    if not session:
        abort(404)
    filename = request.form.get("filename") or ""
    if not filename:
        if _wants_json():
            return _json_error("Filename is required.", status=400)
        flash("Filename is required.", "error")
        return redirect(url_for("view_session", session_id=session_id)), 400
    attachment_path = _session_dir(module_id, session_id) / "attachments" / filename
    if attachment_path.exists():
        attachment_path.unlink()
    _rebuild_attachment_index(_session_dir(module_id, session_id))
    if _wants_json():
        return jsonify({"ok": True})
    flash("Attachment deleted.", "success")
    return redirect(url_for("view_session", session_id=session_id))


@app.route("/modules/<module_id>/sessions/<session_id>/attachments/<filename>")
def open_attachment(module_id: str, session_id: str, filename: str):
    attachment_path = _session_dir(module_id, session_id) / "attachments" / filename
    if not attachment_path.exists():
        abort(404)
    mime, _ = mimetypes.guess_type(str(attachment_path))
    return send_file(attachment_path, mimetype=mime or "application/octet-stream")


@app.route("/attachments/<attachment_id>/open")
def open_attachment_alias(attachment_id: str):
    session_id = request.args.get("session_id")
    if not session_id:
        abort(404)
    session = db.fetch_one("SELECT * FROM sessions WHERE id = ?", (session_id,))
    if not session:
        abort(404)
    return open_attachment(session["module_id"], session_id, attachment_id)


@app.route("/modules/<module_id>/sessions/<session_id>/attachments/<filename>/preview")
def attachment_preview(module_id: str, session_id: str, filename: str):
    attachment_path = _session_dir(module_id, session_id) / "attachments" / filename
    if not attachment_path.exists():
        abort(404)
    mime, _ = mimetypes.guess_type(str(attachment_path))
    return render_template(
        "attachment_preview.html",
        filename=filename,
        mime=mime or "application/octet-stream",
        open_url=url_for("open_attachment", module_id=module_id, session_id=session_id, filename=filename),
    )


@app.route("/modules/<module_id>/sessions/<session_id>/annotations", methods=["POST"])
def save_annotations(module_id: str, session_id: str):
    session = db.fetch_one(
        "SELECT * FROM sessions WHERE id = ? AND module_id = ?", (session_id, module_id)
    )
    if not session:
        abort(404)
    session_dir = _session_dir(module_id, session_id)
    annotations = _load_annotations(session_dir)
    if request.is_json:
        payload = request.get_json(silent=True) or {}
        tags_list = payload.get("tags") or []
        session_tags = payload.get("session_tags") or []
        notes_html = payload.get("personal_notes_html") or ""
        notes_markdown = payload.get("personal_notes_markdown") or ""
        notes_plain = payload.get("personal_notes") or ""
    else:
        tags_list = request.form.getlist("tags")
        session_tags = request.form.getlist("session_tags")
        notes_html = request.form.get("personal_notes_html", "")
        notes_markdown = request.form.get("personal_notes_markdown", "")
        notes_plain = request.form.get("personal_notes", "")
        if len(session_tags) == 1 and "," in session_tags[0]:
            session_tags = [tag.strip() for tag in session_tags[0].split(",") if tag.strip()]
    if tags_list:
        tags_map: dict[str, list[str]] = {}
        for tag_entry in tags_list:
            if ":" not in tag_entry:
                continue
            segment_id, label = tag_entry.split(":", 1)
            label = label.strip().upper()
            if label not in SEGMENT_TAGS:
                continue
            tags_map.setdefault(segment_id.strip(), []).append(label)
        annotations["tags"] = tags_map
    annotations["notes_html"] = notes_html
    annotations["notes_markdown"] = notes_markdown
    annotations["notes"] = notes_plain
    if session_tags:
        annotations["session_tags"] = session_tags
    _save_annotations(session_dir, annotations)
    if _wants_json():
        return jsonify({"ok": True})
    flash("Annotations saved.", "success")
    return redirect(url_for("view_session", session_id=session_id))


@app.route("/modules/<module_id>/sessions/<session_id>/delete-transcript", methods=["POST"])
def delete_transcript(module_id: str, session_id: str):
    session = db.fetch_one(
        "SELECT * FROM sessions WHERE id = ? AND module_id = ?", (session_id, module_id)
    )
    if not session:
        abort(404)
    _clear_transcript(_session_dir(module_id, session_id))
    return jsonify({"ok": True, "message": "Transcript deleted."})


@app.route("/modules/<module_id>/sessions/<session_id>/segment-tags", methods=["POST"])
def update_segment_tags(module_id: str, session_id: str):
    session = db.fetch_one(
        "SELECT * FROM sessions WHERE id = ? AND module_id = ?", (session_id, module_id)
    )
    if not session:
        abort(404)
    payload = request.get_json(silent=True) or {}
    segment_id = payload.get("segment_id")
    label = (payload.get("label") or "").upper()
    checked = bool(payload.get("checked"))
    if not segment_id or label not in SEGMENT_TAGS:
        return _json_error("Invalid segment tag payload.", status=400)
    session_dir = _session_dir(module_id, session_id)
    annotations = _load_annotations(session_dir)
    tags = annotations.setdefault("tags", {}).get(segment_id, [])
    if checked and label not in tags:
        tags.append(label)
    if not checked and label in tags:
        tags = [tag for tag in tags if tag != label]
    annotations["tags"][segment_id] = tags
    _save_annotations(session_dir, annotations)
    return jsonify({"ok": True, "tags": tags})


@app.route("/modules/<module_id>/sessions/<session_id>/generate-notes", methods=["POST"])
def start_notes(module_id: str, session_id: str):
    session = db.fetch_one(
        "SELECT * FROM sessions WHERE id = ? AND module_id = ?", (session_id, module_id)
    )
    if not session:
        abort(404)
    session_dir = _session_dir(module_id, session_id)
    transcript = load_transcript(session_dir / "transcript" / "transcript.json")
    attachment_sources = _load_extracted_sources(session_dir)
    if not transcript and not attachment_sources:
        if _wants_json():
            return _json_error("Upload transcript or attachments to generate notes.", status=400)
        flash("Upload transcript or attachments to generate notes.", "error")
        return redirect(url_for("view_session", session_id=session_id)), 400
    prompt = _build_notes_prompt(transcript, attachment_sources)

    def _run_notes(progress_cb=None):
        if progress_cb:
            progress_cb(10, "Generating notes...")
        output = generate_notes(prompt)
        notes_dir = session_dir / "notes"
        notes_dir.mkdir(parents=True, exist_ok=True)
        notes_path = notes_dir / "ai_notes.md"
        notes_path.write_text(output.notes_markdown, encoding="utf-8")
        notes_json = notes_dir / "ai_notes.json"
        notes_json.write_text(
            json.dumps({"summary": output.summary, "suggested_tags": output.suggested_tags}, indent=2),
            encoding="utf-8",
        )
        if progress_cb:
            progress_cb(100, "Notes generated.")
        return str(notes_path)

    job_id = create_job("Queued for notes generation.")
    try:
        enqueue_job(job_id, _run_notes)
    except GeminiError as exc:
        if _wants_json():
            return _json_error(exc.user_message, status=500)
        flash(exc.user_message, "error")
        return redirect(url_for("view_session", session_id=session_id)), 500
    if _wants_json():
        return jsonify({"job_id": job_id, "redirect": url_for("view_session", session_id=session_id)})
    flash("Generating notes...", "success")
    return redirect(url_for("view_session", session_id=session_id, job_id=job_id))


@app.route("/modules/<module_id>/sessions/<session_id>/ai-notes", methods=["GET"])
def fetch_ai_notes(module_id: str, session_id: str):
    session = db.fetch_one(
        "SELECT * FROM sessions WHERE id = ? AND module_id = ?", (session_id, module_id)
    )
    if not session:
        abort(404)
    session_dir = _session_dir(module_id, session_id)
    notes_path = session_dir / "notes" / "ai_notes.md"
    if not notes_path.exists():
        return _json_error("AI notes not found.", status=404)
    notes, suggested_tags = _load_ai_notes(session_dir)
    return jsonify({"notes": notes, "suggested_tags": suggested_tags})


@app.route("/modules/<module_id>/sessions/<session_id>/export", methods=["GET", "POST"])
def export_pack(module_id: str, session_id: str):
    session = db.fetch_one(
        "SELECT * FROM sessions WHERE id = ? AND module_id = ?", (session_id, module_id)
    )
    module = db.fetch_one("SELECT * FROM modules WHERE id = ?", (module_id,))
    if not session or not module:
        if _wants_json():
            return _json_error("Session not found.", status=404)
        abort(404)
    if request.method == "GET":
        return render_template("export.html", module=dict(module), session=dict(session))

    def _flag(name: str, default: bool = False) -> bool:
        if name not in request.form:
            return default
        return request.form.get(name) == "1"

    include_ai_notes = _flag("include_ai_notes", True)
    include_personal_notes = _flag("include_personal_notes", True)
    include_transcript = _flag("include_transcript", True)
    include_audio = _flag("include_audio", True)
    include_attachments = _flag("include_attachments", True)
    include_raw_chunks = _flag("include_raw_chunks", False)
    include_prompt_manifest = _flag("include_prompt_manifest", False)

    if not any(
        [
            include_ai_notes,
            include_personal_notes,
            include_transcript,
            include_audio,
            include_attachments,
            include_raw_chunks,
            include_prompt_manifest,
        ]
    ):
        if _wants_json():
            return _json_error("Select at least one item to export.", status=400)
        flash("Select at least one item to export.", "error")
        return redirect(url_for("view_session", session_id=session_id)), 400

    zip_path = build_session_export(
        module=dict(module),
        session=dict(session),
        session_dir=_session_dir(module_id, session_id),
        include_ai_notes=include_ai_notes,
        include_personal_notes=include_personal_notes,
        include_transcript=include_transcript,
        include_audio=include_audio,
        include_attachments=include_attachments,
        include_raw_chunks=include_raw_chunks,
        include_prompt_manifest=include_prompt_manifest,
    )
    return send_file(
        zip_path,
        as_attachment=True,
        download_name=zip_path.name,
        mimetype="application/zip",
    )


@app.route("/modules/<module_id>/sessions/<session_id>/qa", methods=["POST"])
def ask_question(module_id: str, session_id: str):
    session = db.fetch_one(
        "SELECT * FROM sessions WHERE id = ? AND module_id = ?", (session_id, module_id)
    )
    if not session:
        abort(404)
    question = (request.form.get("question") or "").strip()
    scope = (request.form.get("scope") or "session").strip()
    if not question:
        flash("Question is required.", "error")
        return redirect(url_for("view_session", session_id=session_id)), 400
    response = _handle_qa_request(session_id, question, scope)
    if isinstance(response, tuple):
        payload, status = response
        flash(payload.get("error", "Q&A failed."), "error")
        return redirect(url_for("view_session", session_id=session_id)), status
    flash("Answer generated.", "success")
    return redirect(url_for("view_session", session_id=session_id))


@app.route("/api/ai/ask", methods=["POST"])
def api_ai_ask():
    payload = request.get_json(silent=True) or {}
    session_id = payload.get("session_id")
    question = (payload.get("question") or "").strip()
    scope = (payload.get("scope") or "session").strip()
    if not session_id or not question:
        return _json_error("Session and question are required.", status=400)
    session = db.fetch_one("SELECT * FROM sessions WHERE id = ?", (session_id,))
    if not session:
        return _json_error("Session not found.", status=404)
    result = _handle_qa_request(session_id, question, scope)
    if isinstance(result, tuple):
        return result
    return jsonify(result)


@app.route("/api/sessions/<session_id>/ai/messages", methods=["GET"])
def api_ai_messages(session_id: str):
    session = db.fetch_one("SELECT * FROM sessions WHERE id = ?", (session_id,))
    if not session:
        return _json_error("Session not found.", status=404)
    return jsonify({"messages": _load_ai_messages(session_id)})


@app.route("/api/source_preview", methods=["GET"])
@app.route("/api/source-preview", methods=["GET"])
@app.route("/api/sources/<source_id>/preview", methods=["GET"])
def api_source_preview(source_id: str | None = None):
    source_id = source_id or request.args.get("source_id")
    session_id = request.args.get("session_id")
    if not source_id or not session_id:
        return _json_error("source_id and session_id are required.", status=400)
    row = db.fetch_one(
        """
        SELECT s.* FROM ai_message_sources s
        JOIN ai_messages m ON m.id = s.message_id
        WHERE s.source_id = ? AND m.session_id = ?
        ORDER BY s.id DESC
        LIMIT 1
        """,
        (source_id, session_id),
    )
    if not row:
        return _json_error("Source not found.", status=404)
    source_json = {}
    if row["source_json"]:
        try:
            source_json = json.loads(row["source_json"])
        except json.JSONDecodeError:
            source_json = {}
    response = {
        "source_id": row["source_id"],
        "kind": row["kind"],
        "title": source_json.get("title") or row["label"],
        "excerpt": source_json.get("excerpt") or row["snippet"],
        "excerpt_full": source_json.get("excerpt") or row["snippet"],
        "open_url": source_json.get("open_url") or row["url"],
        "meta": source_json.get("locator") or {},
    }
    return jsonify(response)


@app.route("/modules/<module_id>", methods=["PATCH"])
def update_module(module_id: str):
    module = db.fetch_one("SELECT * FROM modules WHERE id = ?", (module_id,))
    if not module:
        return _json_error("Module not found.", status=404)
    payload = request.get_json(silent=True) or {}
    name = (payload.get("name") or "").strip()
    if not name:
        return _json_error("Name is required.", status=400)
    db.execute("UPDATE modules SET name = ? WHERE id = ?", (name, module_id))
    return jsonify({"id": module_id, "name": name})


@app.route("/modules/<module_id>", methods=["DELETE"])
def delete_module(module_id: str):
    module = db.fetch_one("SELECT * FROM modules WHERE id = ?", (module_id,))
    if not module:
        return _json_error("Module not found.", status=404)
    db.execute("DELETE FROM sessions WHERE module_id = ?", (module_id,))
    db.execute("DELETE FROM modules WHERE id = ?", (module_id,))
    module_dir = _module_dir(module_id)
    if module_dir.exists():
        shutil.rmtree(module_dir)
    return jsonify({"redirect": url_for("home")})


@app.route("/sessions/<session_id>", methods=["PATCH"])
def update_session(session_id: str):
    session = db.fetch_one("SELECT * FROM sessions WHERE id = ?", (session_id,))
    if not session:
        return _json_error("Session not found.", status=404)
    payload = request.get_json(silent=True) or {}
    name = (payload.get("name") or "").strip()
    if not name:
        return _json_error("Name is required.", status=400)
    db.execute("UPDATE sessions SET name = ? WHERE id = ?", (name, session_id))
    return jsonify({"id": session_id, "name": name})


@app.route("/sessions/<session_id>", methods=["DELETE"])
def delete_session(session_id: str):
    session = db.fetch_one("SELECT * FROM sessions WHERE id = ?", (session_id,))
    if not session:
        return _json_error("Session not found.", status=404)
    module_id = session["module_id"]
    db.execute("DELETE FROM sessions WHERE id = ?", (session_id,))
    session_dir = _session_dir(module_id, session_id)
    if session_dir.exists():
        shutil.rmtree(session_dir)
    next_session = db.fetch_one(
        "SELECT * FROM sessions WHERE module_id = ? ORDER BY created_at DESC LIMIT 1",
        (module_id,),
    )
    if next_session:
        redirect_url = url_for("view_session", session_id=next_session["id"])
    else:
        redirect_url = url_for("view_module", module_id=module_id)
    return jsonify({"redirect": redirect_url})


@app.errorhandler(TranscriptionError)
def handle_transcription_error(error: TranscriptionError):
    flash(error.user_message or "Transcription failed.", "error")
    return redirect(request.referrer or url_for("home"))


@app.errorhandler(GeminiError)
def handle_gemini_error(error: GeminiError):
    flash(error.user_message or "AI request failed.", "error")
    return redirect(request.referrer or url_for("home"))


def create_app(*, testing: bool = False, data_dir: Path | None = None, db_path: Path | None = None) -> Flask:
    if data_dir or db_path:
        config.override_paths(data_dir=data_dir, db_path=db_path)
    app.config["TESTING"] = testing
    _init()
    return app


_init()
